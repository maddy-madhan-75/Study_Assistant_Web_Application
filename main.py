# main.py
import os
import json
import streamlit as st
from dotenv import load_dotenv
from typing import Optional, Any

# ---------- Libraries for file extraction ----------
import PyPDF2
import docx
import openpyxl
from pptx import Presentation
import requests
from bs4 import BeautifulSoup

# ---------- Gemini client ----------
from google import genai
from google.genai import types

# ---------- Structured output types (your dataclasses) ----------
try:
    from lib.quizzing_engine import Quiz, FlashcardSet
except Exception:
    Quiz = None
    FlashcardSet = None

# ---------- Load environment ----------
load_dotenv()
API_KEY = os.getenv("GEMINI_API_KEY")

if not API_KEY:
    st.error("❌ GEMINI_API_KEY not set. Add it to a .env file or Windows environment.")
    st.stop()

# ---------- Initialize Gemini ----------
try:
    client = genai.Client(api_key=API_KEY)
    MODEL = "gemini-2.5-flash"
except Exception as e:
    st.error(f"❌ Failed to initialize Gemini: {e}")
    st.stop()

# ---------------------------------------------------------
# Gemini response extractors
# ---------------------------------------------------------
def extract_text_from_response(response: Any) -> str:
    try:
        return response.candidates[0].content.parts[0].text
    except Exception:
        pass
    try:
        if hasattr(response, "output"):
            return str(response.output)
    except Exception:
        pass
    return str(response)

def extract_structured_from_response(response: Any) -> Any:
    try:
        if hasattr(response, "output"):
            return response.output
    except Exception:
        pass
    try:
        cand = response.candidates[0]
        text = cand.content.parts[0].text
        return json.loads(text)
    except Exception:
        return None

# ---------------------------------------------------------
# File extraction handlers
# ---------------------------------------------------------
def extract_text_from_pdf(file):
    file.seek(0)
    reader = PyPDF2.PdfReader(file)
    pages = []
    for p in reader.pages:
        pages.append(p.extract_text() or "")
    return "\n".join(pages)

def extract_text_from_docx(file):
    file.seek(0)
    document = docx.Document(file)
    return "\n".join([p.text for p in document.paragraphs])

def extract_text_from_excel(file):
    file.seek(0)
    wb = openpyxl.load_workbook(file, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join(str(cell) for cell in row if cell)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)

def extract_text_from_pptx(file):
    """
    Extract text from a .pptx file robustly.
    Uses shape.has_text_frame / shape.text_frame.text when available,
    falls back to getattr(shape, 'text', None) for compatibility.
    """
    file.seek(0)
    prs = Presentation(file)
    parts = []

    for slide in prs.slides:
        for shape in slide.shapes:
            # Preferred: use has_text_frame which is explicit in python-pptx
            try:
                has_tf = getattr(shape, "has_text_frame", False)
            except Exception:
                has_tf = False

            if has_tf:
                # text_frame may contain several paragraphs; .text gives joined text
                try:
                    tf = getattr(shape, "text_frame", None)
                    if tf is not None:
                        txt = getattr(tf, "text", None)
                        if txt and txt.strip():
                            parts.append(txt)
                        continue
                except Exception:
                    pass

            # Fallback: some shapes expose .text directly (rare); use getattr safely
            txt = getattr(shape, "text", None)
            if txt and isinstance(txt, str) and txt.strip():
                parts.append(txt)

    return "\n".join(parts)


def extract_text_from_url(url):
    resp = requests.get(url, timeout=10)
    resp.raise_for_status()
    soup = BeautifulSoup(resp.text, "html.parser")
    for s in soup(["script", "style", "noscript"]):
        s.extract()
    return soup.get_text(separator="\n")

def extract_text_router(uploaded_file, url_input) -> Optional[str]:
    if url_input:
        try:
            return extract_text_from_url(url_input)
        except Exception as e:
            st.sidebar.error(f"Error loading URL: {e}")
            return None

    if uploaded_file:
        name = uploaded_file.name.lower()
        try:
            if name.endswith(".pdf"):
                return extract_text_from_pdf(uploaded_file)
            if name.endswith(".docx"):
                return extract_text_from_docx(uploaded_file)
            if name.endswith(".xlsx"):
                return extract_text_from_excel(uploaded_file)
            if name.endswith(".pptx"):
                return extract_text_from_pptx(uploaded_file)
            if name.endswith(".txt"):
                uploaded_file.seek(0)
                return uploaded_file.read().decode("utf-8")
        except Exception as e:
            st.sidebar.error(f"Error extracting text: {e}")
            return None

    return None

# ---------------------------------------------------------
# Gemini generation functions
# ---------------------------------------------------------
@st.cache_data(show_spinner=False)
def generate_summary(content: str) -> str:
    prompt = f"""
You are an expert study assistant. Provide a concise, structured markdown summary.

STUDY MATERIAL:
---
{content}
---
"""
    try:
        resp = client.models.generate_content(model=MODEL, contents=prompt)
        return extract_text_from_response(resp)
    except Exception as e:
        return f"Error generating summary: {e}"

@st.cache_data(show_spinner=False)
def generate_quiz(content: str):
    prompt = f"""
You are an expert test creator. Create 5 high-quality MCQs, each with 4 options.
Return JSON or structured output.

STUDY MATERIAL:
---
{content}
---
"""
    try:
        cfg = None
        if Quiz:
            try:
                cfg = types.GenerateContentConfig(
                    response_mime_type="application/json",
                    response_schema=Quiz
                )
            except Exception:
                cfg = None

        if cfg:
            resp = client.models.generate_content(
                model=MODEL, contents=prompt, config=cfg)
            return extract_structured_from_response(resp)

        # fallback JSON
        resp = client.models.generate_content(model=MODEL, contents=prompt)
        text = extract_text_from_response(resp)
        try:
            return json.loads(text)
        except Exception:
            return text
    except Exception as e:
        st.error(f"Quiz error: {e}")
        return None

@st.cache_data(show_spinner=False)
def generate_flashcards(content: str):
    prompt = f"""
You are a flashcard expert. Create 10 key-term flashcards as JSON with 'front' and 'back'.

STUDY MATERIAL:
---
{content}
---
"""
    try:
        cfg = None
        if FlashcardSet:
            try:
                cfg = types.GenerateContentConfig(
                    response_mime_type="application/json",
                    response_schema=FlashcardSet
                )
            except Exception:
                cfg = None

        if cfg:
            resp = client.models.generate_content(
                model=MODEL, contents=prompt, config=cfg)
            return extract_structured_from_response(resp)

        # fallback JSON
        resp = client.models.generate_content(model=MODEL, contents=prompt)
        text = extract_text_from_response(resp)
        try:
            return json.loads(text)
        except Exception:
            return text
    except Exception as e:
        st.error(f"Flashcard error: {e}")
        return None

@st.cache_data(show_spinner=False)
def generate_study_plan(content: str) -> str:
    prompt = f"""
You are a professional study strategist. Create a structured 7-day study plan.

STUDY MATERIAL:
---
{content}
---
"""
    try:
        resp = client.models.generate_content(model=MODEL, contents=prompt)
        return extract_text_from_response(resp)
    except Exception as e:
        return f"Error generating study plan: {e}"

# ---------------------------------------------------------
# Streamlit App
# ---------------------------------------------------------
st.set_page_config(page_title="Study Assistant Agent", layout="wide")
st.title("📚 AI Study Assistant")
st.caption(f"Powered by {MODEL}")

# Upload widgets
uploaded_file = st.sidebar.file_uploader(
    "Upload Study Material (TXT, PDF, DOCX, XLSX, PPTX)",
    type=["txt", "pdf", "docx", "xlsx", "pptx"]
)
url_input = st.sidebar.text_input("Or paste a website URL")

# Extract text
study_content = None
if uploaded_file or url_input:
    with st.spinner("Extracting text..."):
        study_content = extract_text_router(uploaded_file, url_input)
        if not study_content:
            st.stop()

st.session_state.file_content = study_content or st.session_state.get("file_content")

if st.session_state.file_content:
    study_content = st.session_state.file_content

    st.subheader("Content Preview")
    st.text_area("Extracted content", study_content[:500] + "..." if len(study_content) > 500 else study_content,
                 height=150, disabled=True)

    tab1, tab2, tab3, tab4 = st.tabs(["📝 Summary", "❓ Quiz", "🧠 Flashcards", "📅 Study Plan"])

    # Summary tab
    with tab1:
        if st.button("Generate Summary"):
            with st.spinner("Generating summary..."):
                st.markdown(generate_summary(study_content))

    # Quiz tab
    with tab2:
        if st.button("Generate Quiz"):
            with st.spinner("Generating quiz..."):
                quiz_data = generate_quiz(study_content)
                st.session_state.quiz_data = quiz_data
                st.session_state.graded = False
                st.rerun()

        if st.session_state.get("quiz_data"):
            quiz = st.session_state.quiz_data
            # handle dict or list
            if isinstance(quiz, dict):
                questions = quiz.get("questions", quiz)
                title = quiz.get("quiz_title", "Generated Quiz")
            else:
                questions = quiz
                title = "Generated Quiz"

            st.subheader(f"Quiz: {title}")

            total_correct = 0
            try:
                for i, q in enumerate(questions):
                    # normalize q
                    if isinstance(q, dict):
                        question_text = q.get("question") or str(q)
                        options = q.get("options", [])
                        correct = q.get("correct_answer")
                        explanation = q.get("explanation", "")
                    else:
                        question_text = getattr(q, "question", str(q))
                        options = getattr(q, "options", [])
                        correct = getattr(q, "correct_answer", None)
                        explanation = getattr(q, "explanation", "")

                    st.write(f"### Q{i+1}: {question_text}")
                    key = f"q{i}"
                    user_choice = st.radio("Choose:", options, key=key)

                    if st.session_state.get("graded"):
                        if user_choice == correct:
                            st.success("Correct")
                            total_correct += 1
                        else:
                            st.error("Incorrect")
                        st.info(f"Correct: {correct}\n\n{explanation}")
            except Exception as e:
                st.error(f"Error rendering quiz: {e}")

            if not st.session_state.get("graded"):
                if st.button("Submit Answers"):
                    st.session_state.graded = True
                    st.rerun()
            else:
                count_qs = len(list(questions)) if questions else 0
                st.metric("Score", f"{total_correct}/{count_qs}")

    # --- TAB 3: FLASHCARDS (robust display) ---
    with tab3:
        st.subheader("Key Term Flashcards")

        if st.button("Generate 10 Flashcards", key="btn_generate_cards"):
            with st.spinner("Creating flashcard deck..."):
                st.session_state.flashcards = generate_flashcards(study_content)

        raw_cards = st.session_state.get("flashcards")
        if raw_cards:
            # unwrap wrapper shapes like {"flashcards": [...]}
            if isinstance(raw_cards, dict) and "flashcards" in raw_cards:
                cards = raw_cards["flashcards"]
            else:
                cards = raw_cards

            # If a single JSON string was returned, try to parse it
            if isinstance(cards, str):
                try:
                    parsed = json.loads(cards)
                    cards = parsed
                except Exception:
                    cards = [cards]

            # Normalize to iterable
            if not isinstance(cards, (list, tuple)):
                cards = [cards]

            for i, card in enumerate(cards):
                front = None
                back = None

                if isinstance(card, dict):
                    front = card.get("front") or card.get("term") or card.get("question")
                    back = card.get("back") or card.get("definition") or card.get("answer")
                else:
                    front = getattr(card, "front", None) or getattr(card, "term", None) or getattr(card, "question", None)
                    back = getattr(card, "back", None) or getattr(card, "definition", None) or getattr(card, "answer", None)

                if front is None and isinstance(card, (list, tuple)) and len(card) > 0:
                    front = card[0]
                    back = card[1] if len(card) > 1 else ""

                if front is None:
                    try:
                        front = str(card)
                    except Exception:
                        front = "<unknown card>"

                if back is None:
                    try:
                        back = str(getattr(card, "back", "")) or ""
                    except Exception:
                        back = ""

                front = str(front)
                back = str(back)

                with st.expander(f"Card {i+1}: {front}"):
                    if back.strip():
                        st.write(back)
                    else:
                        st.write("_No definition provided._")

            # Optional debug
            if st.sidebar.checkbox("Show raw flashcards (debug)", value=False):
                st.sidebar.write(raw_cards)

    # --- TAB 4: STUDY PLAN (robust display) ---
    with tab4:
        st.subheader("7-Day Study Strategy")
        st.markdown("Generate a week-long plan to systematically cover the material.")

        if st.button("Generate Study Plan", key="btn_generate_plan"):
            with st.spinner("Analyzing material and creating schedule..."):
                plan_raw = generate_study_plan(study_content)

                # Normalize plan output
                plan_text = ""
                if isinstance(plan_raw, str):
                    plan_text = plan_raw
                else:
                    try:
                        plan_text = json.dumps(plan_raw, indent=2, ensure_ascii=False)
                    except Exception:
                        try:
                            plan_text = str(plan_raw)
                        except Exception:
                            plan_text = "Unable to render study plan."

                # Display
                # If looks like markdown or multiline, render markdown; otherwise show as code
                if isinstance(plan_text, str) and (plan_text.strip().startswith("#") or "\n" in plan_text):
                    st.markdown(plan_text)
                else:
                    st.code(plan_text)

else:
    st.info("Upload a file or enter a URL to begin.")
